import io
import fitz  # PyMuPDF
from flask import Blueprint, render_template, request, send_file, jsonify
from PIL import Image
import img2pdf
from docx import Document as DocxDocument
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT

try:
    from pdf2docx import Converter as Pdf2DocxConverter
    HAS_PDF2DOCX = True
except ImportError:
    HAS_PDF2DOCX = False

# Marker is loaded lazily inside the route to avoid the ~2GB model preload
# on server start. We only check importability here.
try:
    import marker  # type: ignore
    HAS_MARKER = True
except ImportError:
    HAS_MARKER = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

try:
    import ezdxf
    from ezdxf.addons.drawing import RenderContext, Frontend
    from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    HAS_EZDXF = True
except ImportError:
    HAS_EZDXF = False

import shutil
from routes._helpers import safe_int, safe_float, log_error, NO_FILE_SINGLE, NO_FILE_MULTIPLE

ODA_CONVERTER = shutil.which("ODAFileConverter") or shutil.which("oda_file_converter")
SOFFICE = shutil.which("soffice") or shutil.which("libreoffice")

try:
    from pptx import Presentation
    from pptx.util import Emu
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

bp = Blueprint("convert", __name__)


# ── LibreOffice availability note (PPT/ODP/DOC conversion) ──────

def _soffice_available_notes():
    if SOFFICE:
        return (
            f'<p><i class="bi bi-check-circle-fill" style="color:#2ec4b6"></i> '
            f'<strong>LibreOffice detected:</strong> <code>{SOFFICE}</code></p>'
        )
    return (
        '<p><i class="bi bi-exclamation-triangle-fill" style="color:#ffb703"></i> '
        '<strong>LibreOffice was not found on PATH.</strong> '
        'This tool will not work until LibreOffice is installed.</p>'
        '<details><summary>How to install LibreOffice</summary>'
        '<p><strong>Windows:</strong> Download from '
        '<a href="https://www.libreoffice.org/download/download-libreoffice/" target="_blank">libreoffice.org</a> '
        'and add the install’s <code>program</code> folder '
        '(usually <code>C:\\Program Files\\LibreOffice\\program</code>) to your PATH, '
        'then restart the server.</p>'
        '<p><strong>macOS:</strong> <code>brew install --cask libreoffice</code> '
        '(the <code>soffice</code> binary lives at '
        '<code>/Applications/LibreOffice.app/Contents/MacOS/soffice</code>).</p>'
        '<p><strong>Linux:</strong> <code>sudo apt install libreoffice</code> '
        '(Debian/Ubuntu) or <code>sudo dnf install libreoffice</code> (Fedora).</p>'
        '<p>Restart the server after installing so the new PATH is picked up.</p>'
        '</details>'
    )


def _soffice_convert(file_data: bytes, source_ext: str, target_ext: str = "pdf",
                     timeout: int = 180):
    """Run LibreOffice's headless converter on the given bytes.

    Returns the converted file as bytes on success, or None if soffice is not
    available / the conversion failed (caller falls back to a different engine).

    `source_ext` is used for the temp filename (e.g. "docx", "html", "xlsx").
    """
    if not SOFFICE:
        return None
    import os
    import subprocess
    import tempfile

    with tempfile.TemporaryDirectory() as tmp:
        in_path = os.path.join(tmp, f"input.{source_ext}")
        with open(in_path, "wb") as fp:
            fp.write(file_data)
        try:
            proc = subprocess.run(
                [SOFFICE, "--headless", "--convert-to", target_ext,
                 "--outdir", tmp, in_path],
                capture_output=True, timeout=timeout,
            )
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            log_error(e, f"soffice {source_ext}->{target_ext}")
            return None
        if proc.returncode != 0:
            err = proc.stderr.decode("utf-8", errors="replace")[:200]
            log_error(RuntimeError(err), f"soffice {source_ext}->{target_ext}")
            return None
        out_path = os.path.join(tmp, f"input.{target_ext}")
        if not os.path.exists(out_path):
            return None
        with open(out_path, "rb") as fp:
            return fp.read()


# ── Page Routes ──────────────────────────────────

@bp.route("/to-pdf")
def to_pdf_page():
    if SOFFICE:
        notes = (
            f'<p><i class="bi bi-check-circle-fill" style="color:#2ec4b6"></i> '
            f'<strong>LibreOffice detected</strong> — Word documents will be converted with full layout fidelity '
            f'(fonts, tables, columns, headers/footers preserved).</p>'
        )
    else:
        notes = (
            '<p><i class="bi bi-info-circle-fill" style="color:#4361ee"></i> '
            '<strong>Tip:</strong> install LibreOffice for much better Word→PDF layout fidelity. '
            'Without it, Word files are converted with a basic reflow that loses styling. '
            'See the <a href="/convert/pptx-to-pdf">PowerPoint to PDF</a> page for install instructions.</p>'
        )
    return render_template("upload_tool.html",
        title="Files to PDF",
        description="Convert images, Word documents, and text files to PDF",
        notes=notes,
        endpoint="/convert/to-pdf",
        accept=".jpg,.jpeg,.png,.bmp,.tiff,.webp,.txt,.docx,.doc,.odt",
        multiple=True,
        options=[])


@bp.route("/pdf-to-word")
def pdf_to_word_page():
    marker_status = (
        '<li><strong>Marker (ML)</strong> — uses an ML model for structure understanding. '
        'Best fidelity for academic papers, books, and complex documents. <em>'
        + ('Detected and ready.' if HAS_MARKER else 'Not installed — run <code>pip install marker-pdf</code>. '
           'First run downloads ~2 GB of models. Conversion is slow on CPU (30–60s/page).')
        + '</em></li>'
    )
    return render_template("upload_tool.html",
        title="PDF to Word",
        description="Convert PDF documents to Word (.docx) format",
        notes=(
            '<p><strong>Four conversion modes — pick the one that fits your document:</strong></p>'
            '<ul style="margin:.4rem 0 .6rem 1.2rem">'
            '<li><strong>Layout (default)</strong> — uses <code>pdf2docx</code> to preserve tables, columns, and figure positions. '
            'Best for forms, reports, invoices.</li>'
            '<li><strong>Smart structure</strong> — analyses font sizes to detect headings, lists, and paragraphs, '
            'and emits a Word doc with proper heading styles (visible in Word\'s navigation pane). '
            'Best for articles, blog posts, books, and documentation. Drops tables and figures.</li>'
            '<li><strong>Flowing text</strong> — extracts text in reading order, emits one paragraph per block. '
            'No structure detection. Always produces clean output even on awkward PDFs.</li>'
            f'{marker_status}'
            '</ul>'
            '<p style="font-size:.9em;color:var(--muted)">If your PDF is a scan, run it through <a href="/convert/ocr-pdf">OCR PDF</a> first.</p>'
        ),
        endpoint="/convert/pdf-to-word",
        accept=".pdf",
        multiple=False,
        options=[
            {"type": "select", "name": "mode", "label": "Mode", "default": "layout",
             "choices": [
                 {"value": "layout",    "label": "Layout — preserve tables, columns, figures"},
                 {"value": "structure", "label": "Smart structure — detect headings & lists"},
                 {"value": "text",      "label": "Flowing text — clean paragraphs, no structure"},
                 {"value": "marker",    "label": "Marker (ML) — best fidelity, slow, needs install"},
             ]},
            {"type": "text", "name": "pages", "label": "Pages (blank = all)",
             "placeholder": "e.g. 1-3, 5, 8-10"},
            {"type": "checkbox", "name": "extract_tables",
             "label": "Layout mode: detect borderless tables",
             "check_label": "Try harder to find tables (slower, sometimes invents tables)",
             "default": False},
        ])


@bp.route("/pdf-to-images")
def pdf_to_images_page():
    return render_template("upload_tool.html",
        title="PDF to Images",
        description="Convert each PDF page to an image",
        endpoint="/convert/pdf-to-images",
        accept=".pdf",
        multiple=False,
        options=[
            {"type": "select", "name": "format", "label": "Image Format",
             "choices": [
                 {"value": "png", "label": "PNG"},
                 {"value": "jpg", "label": "JPG"},
             ]},
            {"type": "number", "name": "dpi", "label": "Resolution (DPI)", "default": 200, "min": 72, "max": 600},
        ])


@bp.route("/pdf-to-text")
def pdf_to_text_page():
    return render_template("upload_tool.html",
        title="PDF to Text",
        description="Extract all text content from a PDF document",
        endpoint="/convert/pdf-to-text",
        accept=".pdf",
        multiple=False,
        options=[])


@bp.route("/md-to-pdf")
def md_to_pdf_page():
    return render_template("tools/md_to_pdf.html")


@bp.route("/md-to-docx")
def md_to_docx_page():
    return render_template("tools/md_to_docx.html")


@bp.route("/pdf-to-excel")
def pdf_to_excel_page():
    return render_template("upload_tool.html",
        title="PDF to Excel",
        description="Extract tables from a PDF into an .xlsx workbook",
        notes=(
            "<p><strong>How table detection works:</strong> we try both detection strategies in "
            "order of accuracy:</p>"
            "<ul style='margin:.4rem 0 .6rem 1.2rem'>"
            "<li><strong>Auto (recommended)</strong> — tries ruled-line detection first; if a "
            "page has no visible table borders, falls back to text-alignment detection (catches "
            "borderless tables in financial reports, invoices, schedules).</li>"
            "<li><strong>Lines only</strong> — only tables with visible borders. Most accurate "
            "but misses borderless tables.</li>"
            "<li><strong>Text alignment only</strong> — finds tables by detecting columns of "
            "aligned text. Catches borderless tables but can occasionally false-positive on "
            "multi-column body text.</li>"
            "</ul>"
            "<p style='font-size:.9em;color:var(--muted)'><strong>Still get \"no tables found\"?</strong> "
            "Try our <a href='/convert/pdf-to-word'>PDF to Word</a> tool in <em>Layout</em> mode "
            "instead — it uses <code>pdf2docx</code> which is more aggressive about table "
            "detection. If your PDF is scanned, run it through "
            "<a href='/convert/ocr-pdf'>OCR PDF</a> first.</p>"
        ),
        endpoint="/convert/pdf-to-excel",
        accept=".pdf",
        multiple=False,
        options=[
            {"type": "text", "name": "pages", "label": "Pages (leave empty for all)",
             "placeholder": "e.g. 1-3, 5"},
            {"type": "select", "name": "strategy", "label": "Table detection strategy", "default": "auto",
             "choices": [
                 {"value": "auto",  "label": "Auto — lines first, fall back to text alignment"},
                 {"value": "lines", "label": "Lines only (ruled tables)"},
                 {"value": "text",  "label": "Text alignment only (borderless tables)"},
             ]},
            {"type": "select", "name": "mode", "label": "Extraction mode", "default": "tables",
             "choices": [
                 {"value": "tables", "label": "Tables only (recommended)"},
                 {"value": "tables_text", "label": "Tables, fall back to text rows when none are found"},
                 {"value": "text", "label": "Text only — every line becomes a row"},
             ]},
            {"type": "select", "name": "organize", "label": "Sheet organization", "default": "per_table",
             "choices": [
                 {"value": "per_table", "label": "One sheet per detected table"},
                 {"value": "per_page", "label": "One sheet per page (tables stacked)"},
                 {"value": "combined", "label": "Everything on one sheet"},
             ]},
        ])


OCR_LANGS = [
    {"value": "eng", "label": "English"},
    {"value": "ind", "label": "Indonesian"},
    {"value": "fra", "label": "French"},
    {"value": "deu", "label": "German"},
    {"value": "spa", "label": "Spanish"},
    {"value": "ita", "label": "Italian"},
    {"value": "por", "label": "Portuguese"},
    {"value": "rus", "label": "Russian"},
    {"value": "chi_sim", "label": "Chinese (Simplified)"},
    {"value": "chi_tra", "label": "Chinese (Traditional)"},
    {"value": "jpn", "label": "Japanese"},
    {"value": "kor", "label": "Korean"},
    {"value": "ara", "label": "Arabic"},
    {"value": "hin", "label": "Hindi"},
]


@bp.route("/ocr-pdf")
def ocr_pdf_page():
    return render_template("upload_tool.html",
        title="OCR PDF",
        description="Extract text from scanned PDFs or create a searchable PDF with a hidden text layer",
        endpoint="/convert/ocr-pdf",
        accept=".pdf",
        multiple=False,
        options=[
            {"type": "select", "name": "mode", "label": "Output",
             "choices": [
                 {"value": "searchable", "label": "Searchable PDF (image + text layer)"},
                 {"value": "text", "label": "Extracted text only"},
             ]},
            {"type": "select", "name": "lang", "label": "Language",
             "choices": OCR_LANGS},
            {"type": "number", "name": "dpi", "label": "OCR Resolution (DPI)",
             "default": 200, "min": 100, "max": 400},
        ])


@bp.route("/cad-to-pdf")
def cad_to_pdf_page():
    if ODA_CONVERTER:
        notes = (
            '<p><i class="bi bi-check-circle-fill" style="color:#2ec4b6"></i> '
            '<strong>DWG support is enabled.</strong> ODA File Converter was detected at '
            f'<code>{ODA_CONVERTER}</code>.</p>'
            '<p>DXF files are rendered directly. DWG files are auto-converted to DXF first.</p>'
        )
    else:
        notes = (
            '<p><strong>DXF works out of the box.</strong> DWG files need the free '
            '<a href="https://www.opendesign.com/guestfiles/oda_file_converter" target="_blank" rel="noopener">'
            'ODA File Converter</a> installed and available on your system <code>PATH</code>.</p>'
            '<details>'
            '<summary>How to install ODA File Converter</summary>'
            '<ol>'
            '<li>Download the installer for your OS from '
            '<a href="https://www.opendesign.com/guestfiles/oda_file_converter" target="_blank" rel="noopener">opendesign.com</a> '
            '(free, guest download — no account required).</li>'
            '<li>Run the installer. Defaults are fine.</li>'
            '<li><strong>Add it to your PATH so this app can find it:</strong>'
            '<ul>'
            '<li><strong>Windows:</strong> add <code>C:\\Program Files\\ODA\\ODAFileConverter_title_version</code> '
            '(the folder containing <code>ODAFileConverter.exe</code>) to your <em>System Environment Variables</em> &rarr; <code>Path</code>.</li>'
            '<li><strong>macOS:</strong> <code>ln -s /Applications/ODAFileConverter.app/Contents/MacOS/ODAFileConverter /usr/local/bin/ODAFileConverter</code></li>'
            '<li><strong>Linux:</strong> the <code>.deb</code>/<code>.rpm</code> package installs <code>ODAFileConverter</code> on PATH automatically. Otherwise symlink the binary into <code>/usr/local/bin</code>.</li>'
            '</ul></li>'
            '<li>Open a new terminal and verify: <code>ODAFileConverter</code> (should launch the tool GUI, or exit silently).</li>'
            '<li><strong>Restart this Flask server</strong> so it picks up the updated PATH.</li>'
            '</ol>'
            '<p style="margin-top:.4rem">Alternative: open your DWG in free tools like <a href="https://www.autodesk.com/viewers" target="_blank" rel="noopener">Autodesk Viewer</a>, LibreCAD, or QCAD and export it as DXF, then upload the DXF here.</p>'
            '</details>'
        )

    return render_template("upload_tool.html",
        title="CAD to PDF/Image",
        description="Convert DXF drawings to PDF or PNG. DWG is supported when ODA File Converter is installed.",
        notes=notes,
        endpoint="/convert/cad-to-pdf",
        accept=".dxf,.dwg",
        multiple=False,
        options=[
            {"type": "select", "name": "format", "label": "Output Format",
             "choices": [
                 {"value": "pdf", "label": "PDF"},
                 {"value": "png", "label": "PNG"},
             ]},
            {"type": "number", "name": "dpi", "label": "PNG Resolution (DPI)",
             "default": 150, "min": 72, "max": 600,
             "depends_on": {"format": "png"}},
        ])


@bp.route("/html-to-pdf")
def html_to_pdf_page():
    if SOFFICE:
        notes = (
            f'<p><i class="bi bi-check-circle-fill" style="color:#2ec4b6"></i> '
            f'<strong>LibreOffice detected</strong> — full CSS support, tables, lists, and inline styles render correctly.</p>'
        )
    else:
        notes = (
            '<p><i class="bi bi-info-circle-fill" style="color:#4361ee"></i> '
            '<strong>Tip:</strong> install LibreOffice for far better CSS / table / image fidelity. '
            'Without it, PDF rendering uses PyMuPDF\'s minimal HTML parser (basic text and simple tables only).</p>'
        )
    return render_template("upload_tool.html",
        title="HTML to PDF",
        description="Convert HTML content to a PDF document",
        notes=notes,
        endpoint="/convert/html-to-pdf",
        text_input=True,
        text_label="HTML Content",
        text_placeholder="<h1>Hello World</h1>\n<p>Paste your HTML here...</p>",
        accept="",
        multiple=False,
        options=[],
        button_text="Convert to PDF")


# ── Helpers ──────────────────────────────────────

def _docx_to_pdf(data: bytes) -> bytes:
    """Convert a .docx file (as bytes) to PDF bytes using python-docx + reportlab."""
    doc = DocxDocument(io.BytesIO(data))
    buf = io.BytesIO()

    styles = getSampleStyleSheet()
    normal = styles["Normal"]
    normal.fontName = "Helvetica"
    normal.fontSize = 11
    normal.leading = 14

    heading_styles = {}
    for level in range(1, 4):
        size = {1: 18, 2: 15, 3: 13}[level]
        heading_styles[level] = ParagraphStyle(
            f"Heading{level}", parent=normal,
            fontName="Helvetica-Bold", fontSize=size, leading=size + 4,
            spaceBefore=12, spaceAfter=6,
        )

    pdf = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    story = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 6))
            continue

        style_name = para.style.name.lower() if para.style else ""

        if "heading 1" in style_name:
            story.append(Paragraph(text, heading_styles[1]))
        elif "heading 2" in style_name:
            story.append(Paragraph(text, heading_styles[2]))
        elif "heading 3" in style_name:
            story.append(Paragraph(text, heading_styles[3]))
        else:
            # Preserve basic inline formatting
            rich = _build_rich_text(para)
            story.append(Paragraph(rich, normal))

    # Handle tables
    for table in doc.tables:
        tdata = []
        for row in table.rows:
            tdata.append([cell.text for cell in row.cells])
        if tdata:
            t = Table(tdata, repeatRows=1)
            t.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.9, 0.9, 0.95)),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ]))
            story.append(Spacer(1, 8))
            story.append(t)
            story.append(Spacer(1, 8))

    if not story:
        story.append(Paragraph("(empty document)", normal))

    pdf.build(story)
    return buf.getvalue()


def _build_rich_text(para) -> str:
    """Convert a python-docx paragraph's runs into reportlab-compatible rich text."""
    parts = []
    for run in para.runs:
        text = run.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if not text:
            continue
        if run.bold and run.italic:
            parts.append(f"<b><i>{text}</i></b>")
        elif run.bold:
            parts.append(f"<b>{text}</b>")
        elif run.italic:
            parts.append(f"<i>{text}</i>")
        elif run.underline:
            parts.append(f"<u>{text}</u>")
        else:
            parts.append(text)
    return "".join(parts) or para.text


# ── Processing Routes ────────────────────────────

@bp.route("/to-pdf", methods=["POST"])
def to_pdf():
    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_MULTIPLE), 400

    pdf_doc = fitz.open()

    for f in files:
        name = f.filename.lower()
        data = f.read()

        if name.endswith((".docx", ".doc", ".odt")):
            # Word document → PDF pages.
            # Prefer LibreOffice for high-fidelity layout. Fall back to the
            # python-docx + reportlab rebuilder if soffice is unavailable.
            ext = name.rsplit(".", 1)[-1]
            try:
                pdf_bytes = _soffice_convert(data, ext, "pdf")
                if pdf_bytes is None:
                    if ext != "docx":
                        return jsonify(error=(
                            f"'{f.filename}' requires LibreOffice (soffice) on PATH. "
                            "Only .docx is supported by the built-in fallback."
                        )), 400
                    pdf_bytes = _docx_to_pdf(data)
                with fitz.open(stream=pdf_bytes, filetype="pdf") as docx_pdf:
                    pdf_doc.insert_pdf(docx_pdf)
            except Exception as e:
                log_error(e, f"to-pdf docx: {f.filename}")
                return jsonify(error=f"Could not convert '{f.filename}' (Word file may be corrupted)."), 400
        elif name.endswith(".txt"):
            # Text file → PDF page
            text = data.decode("utf-8", errors="replace")
            page = pdf_doc.new_page(width=595, height=842)  # A4
            rect = fitz.Rect(50, 50, 545, 792)
            page.insert_textbox(rect, text, fontsize=11, fontname="helv")
        else:
            # Image → PDF page
            try:
                with Image.open(io.BytesIO(data)) as pil_img:
                    if pil_img.mode in ("RGBA", "P"):
                        pil_img = pil_img.convert("RGB")
                    buf = io.BytesIO()
                    pil_img.save(buf, format="JPEG", quality=95)
                img_data = buf.getvalue()

                with fitz.open(stream=img_data, filetype="jpeg") as img_doc:
                    rect = img_doc[0].rect
                    pdf_page = pdf_doc.new_page(width=rect.width, height=rect.height)
                    pdf_page.insert_image(rect, stream=img_data)
            except Exception as e:
                log_error(e, f"to-pdf image: {f.filename}")
                return jsonify(error=f"Could not convert '{f.filename}' (image may be corrupted or unsupported)."), 400

    output = io.BytesIO()
    pdf_doc.save(output)
    pdf_doc.close()
    output.seek(0)

    return send_file(output, mimetype="application/pdf",
                     as_attachment=True, download_name="converted.pdf")


@bp.route("/pdf-to-word", methods=["POST"])
def pdf_to_word():
    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_SINGLE), 400

    mode = request.form.get("mode", "layout")
    pdf_data = files[0].read()
    pages_spec = (request.form.get("pages") or "").strip()
    extract_borderless_tables = request.form.get("extract_tables") == "on"
    base = files[0].filename.rsplit(".", 1)[0]
    docx_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    # Pre-resolve page range against the PDF (used by all modes)
    try:
        with fitz.open(stream=pdf_data, filetype="pdf") as probe:
            total_pages = len(probe)
            if pages_spec:
                from routes.pdf_tools import parse_page_ranges
                target_pages = parse_page_ranges(pages_spec, total_pages)
                if not target_pages:
                    return jsonify(error="No valid pages selected."), 400
            else:
                target_pages = list(range(total_pages))
    except (ValueError, IndexError):
        return jsonify(error="Invalid page range. Use e.g. '1-3, 5, 8-10'."), 400
    except Exception as e:
        log_error(e, "pdf-to-word probe")
        return jsonify(error="Could not open PDF (the file may be corrupted or password-protected)."), 400

    # ── Mode dispatch ──────────────────────────────────────

    if mode == "text":
        try:
            buf = _pdf_to_docx_flowing_text(pdf_data, target_pages)
        except Exception as e:
            log_error(e, "pdf-to-word text")
            return jsonify(error="Could not extract text from the PDF (it may be a scan — try OCR PDF first)."), 400
        return send_file(io.BytesIO(buf), mimetype=docx_mime,
                         as_attachment=True, download_name=f"{base}.docx")

    if mode == "structure":
        try:
            buf = _pdf_to_docx_smart_structure(pdf_data, target_pages)
        except ValueError as e:
            return jsonify(error=str(e)), 400
        except Exception as e:
            log_error(e, "pdf-to-word structure")
            return jsonify(error="Smart-structure analysis failed. Try Flowing text mode instead."), 400
        return send_file(io.BytesIO(buf), mimetype=docx_mime,
                         as_attachment=True, download_name=f"{base}.docx")

    if mode == "marker":
        if not HAS_MARKER:
            return jsonify(error=(
                "Marker mode requires the 'marker-pdf' package. Run: "
                "pip install marker-pdf — first run will download ~2 GB of models."
            )), 400
        try:
            buf = _pdf_to_docx_via_marker(pdf_data, target_pages)
        except Exception as e:
            log_error(e, "pdf-to-word marker")
            return jsonify(error="Marker conversion failed. Check the server log; "
                           "first run downloads ~2 GB and may need extra time."), 400
        return send_file(io.BytesIO(buf), mimetype=docx_mime,
                         as_attachment=True, download_name=f"{base}.docx")

    # ── Layout mode (default) ──────────────────────────────
    if not HAS_PDF2DOCX:
        return jsonify(error="Layout mode requires pdf2docx. Run: pip install pdf2docx — or switch to 'Flowing text' or 'Smart structure' mode."), 400

    import tempfile, os

    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = os.path.join(tmpdir, "input.pdf")
        docx_path = os.path.join(tmpdir, "output.docx")

        with open(pdf_path, "wb") as f:
            f.write(pdf_data)

        # Translate target_pages (set of 0-based) to start/end if contiguous.
        # pdf2docx supports a `pages` list arg directly, which is cleaner.
        cv_kwargs = {"multi_processing": False}
        if pages_spec:
            cv_kwargs["pages"] = target_pages
        if extract_borderless_tables:
            cv_kwargs["parse_stream_table"] = True

        try:
            cv = Pdf2DocxConverter(pdf_path)
            try:
                cv.convert(docx_path, **cv_kwargs)
            finally:
                cv.close()
        except Exception as e:
            log_error(e, "pdf-to-word layout")
            return jsonify(error="Layout conversion failed. Try Smart structure or Flowing text mode instead, or check that the PDF isn't password-protected."), 400

        with open(docx_path, "rb") as f:
            result = io.BytesIO(f.read())

    result.seek(0)
    name = files[0].filename.rsplit(".", 1)[0] + ".docx"
    return send_file(result, mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                     as_attachment=True, download_name=name)


# ── PDF → Word helpers (one per non-pdf2docx mode) ─────────

def _pdf_to_docx_flowing_text(pdf_data: bytes, target_pages: list[int]) -> bytes:
    """Reading-order text extraction → one paragraph per block. No structure."""
    from docx import Document as DocxDocument

    with fitz.open(stream=pdf_data, filetype="pdf") as src:
        doc = DocxDocument()
        for idx, pno in enumerate(target_pages):
            if idx > 0:
                doc.add_page_break()
            page = src[pno]
            blocks = page.get_text("blocks") or []
            blocks.sort(key=lambda b: (round(b[1], 1), round(b[0], 1)))
            for b in blocks:
                text = (b[4] if len(b) > 4 else "").strip()
                if not text:
                    continue
                text = "\n".join(p.strip() for p in text.split("\n") if p.strip())
                for para in text.split("\n\n"):
                    para = para.replace("\n", " ").strip()
                    if para:
                        doc.add_paragraph(para)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()


def _pdf_to_docx_smart_structure(pdf_data: bytes, target_pages: list[int]) -> bytes:
    """Detect headings (by font size), bullet/numbered lists (by line prefix),
    and paragraphs. Emit a .docx with proper Word heading and list styles.

    Drops tables and figures (those need Layout or Marker mode).
    """
    import re
    from collections import Counter
    from docx import Document as DocxDocument

    BULLET_RE = re.compile(r"^[•▪●·\-\*]\s+")
    NUMBER_RE = re.compile(r"^(\d+|[a-zA-Z])[\.\)]\s+")

    with fitz.open(stream=pdf_data, filetype="pdf") as src:
        # Pass 1: collect font sizes to determine the body baseline.
        sizes: list[float] = []
        for pno in target_pages:
            page = src[pno]
            for block in page.get_text("dict")["blocks"]:
                if "lines" not in block:
                    continue
                for line in block["lines"]:
                    for span in line["spans"]:
                        sizes.append(round(span["size"], 1))

        if not sizes:
            raise ValueError("No text found in the selected pages. If the PDF is a scan, run OCR PDF first.")

        body_size = Counter(sizes).most_common(1)[0][0]

        # Pass 2: build the document.
        doc = DocxDocument()
        for idx, pno in enumerate(target_pages):
            if idx > 0:
                doc.add_page_break()
            page = src[pno]
            blocks = [b for b in page.get_text("dict")["blocks"] if "lines" in b]
            # Reading order: top-to-bottom, then left-to-right within tolerance
            blocks.sort(key=lambda b: (round(b["bbox"][1] / 5) * 5, round(b["bbox"][0])))

            for block in blocks:
                lines = []
                spans_meta: list[tuple[float, bool]] = []
                for line in block["lines"]:
                    line_text = "".join(s["text"] for s in line["spans"])
                    if line_text.strip():
                        lines.append(line_text)
                        for s in line["spans"]:
                            # PyMuPDF flag bit 4 (0x10) = bold
                            spans_meta.append((s["size"], bool(s["flags"] & 16)))

                if not lines:
                    continue

                avg_size = sum(s for s, _ in spans_meta) / len(spans_meta)
                bold_ratio = sum(1 for _, b in spans_meta if b) / len(spans_meta)
                full_text = " ".join(line.strip() for line in lines).strip()

                # Heading detection by relative font size
                if avg_size >= body_size * 1.6:
                    doc.add_heading(full_text, level=1)
                elif avg_size >= body_size * 1.3:
                    doc.add_heading(full_text, level=2)
                elif avg_size >= body_size * 1.15 or (
                    avg_size >= body_size * 1.05 and bold_ratio > 0.6 and len(full_text) < 120
                ):
                    doc.add_heading(full_text, level=3)
                # List detection by line prefix
                elif BULLET_RE.match(full_text):
                    doc.add_paragraph(BULLET_RE.sub("", full_text), style="List Bullet")
                elif NUMBER_RE.match(full_text):
                    doc.add_paragraph(NUMBER_RE.sub("", full_text), style="List Number")
                else:
                    doc.add_paragraph(full_text)

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()


def _pdf_to_docx_via_marker(pdf_data: bytes, target_pages: list[int]) -> bytes:
    """Use Marker (ML) to extract structured Markdown, then convert to .docx."""
    import os
    import tempfile
    from marker.converters.pdf import PdfConverter
    from marker.models import create_model_dict
    from marker.output import text_from_rendered

    # If specific pages requested, build a subset PDF first so Marker only
    # processes what's needed (it's slow per page).
    if len(target_pages) != _count_pages(pdf_data):
        pdf_data = _extract_pages(pdf_data, target_pages)

    with tempfile.TemporaryDirectory() as tmp:
        pdf_path = os.path.join(tmp, "input.pdf")
        with open(pdf_path, "wb") as f:
            f.write(pdf_data)

        converter = PdfConverter(artifact_dict=create_model_dict())
        rendered = converter(pdf_path)
        markdown_text, _, _ = text_from_rendered(rendered)

    # Convert the markdown to docx via a reusable HTML→docx walker
    import markdown as md_lib
    html = md_lib.markdown(
        markdown_text,
        extensions=["extra", "sane_lists", "nl2br", "fenced_code", "tables"],
    )
    return _md_html_to_docx_bytes(html)


def _count_pages(pdf_data: bytes) -> int:
    with fitz.open(stream=pdf_data, filetype="pdf") as d:
        return len(d)


def _extract_pages(pdf_data: bytes, page_indices: list[int]) -> bytes:
    """Build a new PDF containing only the listed page indices (0-based)."""
    with fitz.open(stream=pdf_data, filetype="pdf") as src:
        with fitz.open() as out:
            for idx in page_indices:
                out.insert_pdf(src, from_page=idx, to_page=idx)
            buf = io.BytesIO()
            out.save(buf)
            return buf.getvalue()


def _md_html_to_docx_bytes(html: str) -> bytes:
    """Use the same HTML-walking parser md_to_docx uses, but as a callable
    helper so the marker mode can reuse it. Returns docx bytes.
    """
    # We avoid circular imports — import lazily.
    from html.parser import HTMLParser as _HP
    from docx import Document as DocxDocument

    doc = DocxDocument()

    class _P(_HP):
        def __init__(self):
            super().__init__()
            self.cur_para = None
            self.list_stack = []
            self.in_pre = False

        def handle_starttag(self, tag, attrs):
            if tag in ("h1", "h2", "h3", "h4"):
                self.cur_para = doc.add_heading("", level=int(tag[1]))
            elif tag == "p":
                self.cur_para = doc.add_paragraph()
            elif tag == "li":
                style = "List Number" if (self.list_stack and self.list_stack[-1] == "ol") else "List Bullet"
                self.cur_para = doc.add_paragraph(style=style)
            elif tag in ("ul", "ol"):
                self.list_stack.append(tag)
            elif tag in ("strong", "b", "em", "i", "code"):
                pass  # handled in handle_data
            elif tag == "pre":
                self.in_pre = True
                self.cur_para = doc.add_paragraph(style="Intense Quote")
            elif tag == "hr":
                doc.add_paragraph("─" * 40)

        def handle_endtag(self, tag):
            if tag in ("ul", "ol") and self.list_stack:
                self.list_stack.pop()
            if tag == "pre":
                self.in_pre = False
                self.cur_para = None
            if tag in ("h1", "h2", "h3", "h4", "p", "li"):
                self.cur_para = None

        def handle_data(self, data):
            if self.cur_para is None:
                if data.strip():
                    self.cur_para = doc.add_paragraph()
                else:
                    return
            self.cur_para.add_run(data)

    parser = _P()
    parser.feed(html)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@bp.route("/pdf-to-images", methods=["POST"])
def pdf_to_images():
    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_SINGLE), 400

    fmt = request.form.get("format", "png")
    dpi = safe_int(request.form.get("dpi"), 200, min_val=72, max_val=600)

    pdf_data = files[0].read()
    try:
        doc = fitz.open(stream=pdf_data, filetype="pdf")
    except Exception as e:
        log_error(e, "pdf-to-images open")
        return jsonify(error="Could not open PDF (the file may be corrupted or password-protected)."), 400

    from utils.file_utils import make_zip
    images = []
    mat = fitz.Matrix(dpi / 72, dpi / 72)

    try:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            if fmt == "jpg":
                img_bytes = pix.tobytes("jpeg")
                ext = "jpg"
            else:
                img_bytes = pix.tobytes("png")
                ext = "png"
            images.append((f"page_{i + 1}.{ext}", img_bytes))
    finally:
        doc.close()

    if len(images) == 1:
        mime = "image/png" if fmt == "png" else "image/jpeg"
        return send_file(io.BytesIO(images[0][1]), mimetype=mime,
                         as_attachment=True, download_name=images[0][0])

    zip_buf = make_zip(images)
    name = files[0].filename.rsplit(".", 1)[0] + "_images.zip"
    return send_file(zip_buf, mimetype="application/zip",
                     as_attachment=True, download_name=name)


@bp.route("/pdf-to-text", methods=["POST"])
def pdf_to_text():
    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_SINGLE), 400

    pdf_data = files[0].read()
    try:
        doc = fitz.open(stream=pdf_data, filetype="pdf")
    except Exception as e:
        log_error(e, "pdf-to-text open")
        return jsonify(error="Could not open PDF (the file may be corrupted or password-protected)."), 400

    text_parts = []
    try:
        for i, page in enumerate(doc):
            text_parts.append(f"--- Page {i + 1} ---")
            text_parts.append(page.get_text())
    finally:
        doc.close()

    return jsonify(text="\n".join(text_parts))


@bp.route("/pdf-to-excel", methods=["POST"])
def pdf_to_excel():
    import re
    from openpyxl import Workbook
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter
    from routes.pdf_tools import parse_page_ranges

    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_SINGLE), 400

    mode = request.form.get("mode", "tables")
    organize = request.form.get("organize", "per_table")
    strategy = request.form.get("strategy", "auto")
    if strategy not in ("auto", "lines", "text"):
        strategy = "auto"
    pages_spec = request.form.get("pages", "").strip()

    try:
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
    except Exception as e:
        log_error(e, "pdf-to-excel open")
        return jsonify(error="Could not open PDF (the file may be corrupted or password-protected)."), 400

    try:
        target_pages = parse_page_ranges(pages_spec, len(doc))
    except (ValueError, IndexError):
        doc.close()
        return jsonify(error="Invalid page range. Use e.g. '1-3, 5, 8-10'."), 400
    if not target_pages:
        doc.close()
        return jsonify(error="No valid pages selected."), 400

    wb = Workbook()
    wb.remove(wb.active)
    used_names: set[str] = set()
    total_tables = 0
    total_text_pages = 0

    def _safe_name(base: str) -> str:
        name = re.sub(r"[\[\]\*\?\/\\:]", "_", base)[:31] or "Sheet"
        candidate = name
        i = 2
        while candidate in used_names:
            suffix = f"_{i}"
            candidate = (name[: 31 - len(suffix)] + suffix)
            i += 1
        used_names.add(candidate)
        return candidate

    def _write_rows(ws, rows: list[list], start_row: int = 1, header: bool = True) -> int:
        for r_idx, row in enumerate(rows, start=start_row):
            for c_idx, cell in enumerate(row, start=1):
                ws.cell(row=r_idx, column=c_idx, value="" if cell is None else str(cell))
            if header and r_idx == start_row:
                for c_idx in range(1, len(row) + 1):
                    ws.cell(row=r_idx, column=c_idx).font = Font(bold=True)
        return start_row + len(rows)

    def _text_rows(page) -> list[list[str]]:
        lines = page.get_text().splitlines()
        rows = []
        for line in lines:
            line = line.strip()
            if not line:
                continue
            parts = re.split(r"\s{2,}|\t+", line)
            rows.append(parts if parts else [line])
        return rows

    def _find_tables_robust(page) -> list:
        """Detect tables on a page according to the user's chosen strategy.

        PyMuPDF's default `find_tables()` only catches ruled (visible-border)
        tables. Many real-world PDFs use borderless tables where columns are
        aligned by whitespace — those need `strategy="text"`. The "auto" mode
        tries lines first and only falls back to text-based when nothing is
        found, which avoids the false-positive risk of text-detection picking
        up multi-column body text as a "table".
        """
        try:
            if strategy == "lines":
                return list(page.find_tables(strategy="lines"))
            if strategy == "text":
                return list(page.find_tables(
                    strategy="text",
                    vertical_strategy="text",
                    horizontal_strategy="text",
                ))
            # auto: lines, then text fallback
            tables = list(page.find_tables(strategy="lines"))
            if tables:
                return tables
            return list(page.find_tables(
                strategy="text",
                vertical_strategy="text",
                horizontal_strategy="text",
            ))
        except Exception as e:
            log_error(e, f"find_tables strategy={strategy}")
            return []

    # ── "combined" — stream everything into a single sheet ────────────
    if organize == "combined":
        ws = wb.create_sheet(_safe_name("Extracted"))
        next_row = 1
        for pno in target_pages:
            page = doc[pno]
            page_had_content = False

            if mode in ("tables", "tables_text"):
                tables = _find_tables_robust(page)
                for t in tables:
                    rows = t.extract()
                    if not rows:
                        continue
                    ws.cell(row=next_row, column=1, value=f"Page {pno + 1} – table").font = Font(bold=True, italic=True)
                    next_row += 1
                    next_row = _write_rows(ws, rows, start_row=next_row)
                    next_row += 1
                    total_tables += 1
                    page_had_content = True

            if mode == "text" or (mode == "tables_text" and not page_had_content):
                text_rows = _text_rows(page)
                if text_rows:
                    ws.cell(row=next_row, column=1, value=f"Page {pno + 1} – text").font = Font(bold=True, italic=True)
                    next_row += 1
                    next_row = _write_rows(ws, text_rows, start_row=next_row, header=False)
                    next_row += 1
                    total_text_pages += 1

    # ── "per_page" and "per_table" ────────────────────────────────────
    else:
        for pno in target_pages:
            page = doc[pno]
            tables_rows = []  # list of (label, rows)

            if mode in ("tables", "tables_text"):
                for tidx, t in enumerate(_find_tables_robust(page), start=1):
                    rows = t.extract()
                    if rows:
                        tables_rows.append((f"Table {tidx}", rows))
                        total_tables += 1

            if mode == "text" or (mode == "tables_text" and not tables_rows):
                text_rows = _text_rows(page)
                if text_rows:
                    tables_rows.append(("Text", text_rows))
                    total_text_pages += 1

            if not tables_rows:
                continue

            if organize == "per_table":
                for label, rows in tables_rows:
                    is_text = label == "Text"
                    sheet = wb.create_sheet(_safe_name(f"Page{pno + 1}_{label.replace(' ', '')}"))
                    _write_rows(sheet, rows, header=not is_text)
            else:  # per_page
                sheet = wb.create_sheet(_safe_name(f"Page_{pno + 1}"))
                next_row = 1
                for label, rows in tables_rows:
                    is_text = label == "Text"
                    sheet.cell(row=next_row, column=1, value=label).font = Font(bold=True, italic=True)
                    next_row += 1
                    next_row = _write_rows(sheet, rows, start_row=next_row, header=not is_text)
                    next_row += 1

    doc.close()

    if not wb.sheetnames:
        msg = "No tables found on the selected pages."
        if strategy == "lines":
            msg += " Try the 'Text alignment' or 'Auto' strategy — your PDF may use borderless tables."
        elif mode == "tables":
            msg += " Try the 'Tables, fall back to text rows' mode, or use PDF to Word in Layout mode."
        else:
            msg += " If this is a scanned PDF, run it through OCR PDF first; otherwise try PDF to Word in Layout mode."
        return jsonify(error=msg), 400

    # Auto-size columns on every sheet (cap at 60 chars to avoid absurd widths)
    for ws in wb.worksheets:
        for col_idx in range(1, ws.max_column + 1):
            max_len = 0
            for row_idx in range(1, ws.max_row + 1):
                v = ws.cell(row=row_idx, column=col_idx).value
                if v is not None:
                    max_len = max(max_len, len(str(v)))
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 2, 60)

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)

    name = files[0].filename.rsplit(".", 1)[0] + ".xlsx"
    return send_file(
        output,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        as_attachment=True,
        download_name=name,
    )


@bp.route("/md-to-pdf", methods=["POST"])
def md_to_pdf():
    import markdown as md_lib

    md_text = request.form.get("markdown", "").strip()
    if not md_text:
        return jsonify(error="Please enter some Markdown."), 400

    page_size = request.form.get("page_size", "A4").lower()
    font_size = safe_int(request.form.get("font_size"), 11, min_val=8, max_val=18)

    sizes_map = {
        "a4": (595.28, 841.89),
        "letter": (612, 792),
        "legal": (612, 1008),
        "a5": (419.53, 595.28),
    }
    page_w, page_h = sizes_map.get(page_size, sizes_map["a4"])

    html = md_lib.markdown(
        md_text,
        extensions=["extra", "sane_lists", "nl2br", "fenced_code", "tables"],
    )

    margin = 54  # 0.75 inch
    content_rect = fitz.Rect(margin, margin, page_w - margin, page_h - margin)

    css = (
        f"* {{ font-family: sans-serif; font-size: {font_size}pt; line-height: 1.45; }}"
        "h1 { font-size: 1.8em; margin: 0.4em 0 0.3em; }"
        "h2 { font-size: 1.5em; margin: 0.4em 0 0.3em; }"
        "h3 { font-size: 1.2em; margin: 0.4em 0 0.3em; }"
        "p { margin: 0.35em 0; }"
        "ul, ol { margin: 0.3em 0 0.3em 1.2em; }"
        "li { margin: 0.15em 0; }"
        "code { font-family: monospace; background: #f2f2f2; padding: 1px 3px; }"
        "pre { font-family: monospace; background: #f5f5f5; padding: 0.5em; white-space: pre-wrap; }"
        "blockquote { margin: 0.5em 0; padding-left: 0.8em; border-left: 3px solid #bbb; color: #555; }"
        "table { border-collapse: collapse; margin: 0.4em 0; }"
        "th, td { border: 1px solid #999; padding: 0.2em 0.5em; }"
        "hr { border: none; border-top: 1px solid #ccc; margin: 0.6em 0; }"
    )

    # Use PyMuPDF's Story + DocumentWriter for reliable multi-page HTML rendering
    output = io.BytesIO()
    mediabox = fitz.Rect(0, 0, page_w, page_h)
    writer = fitz.DocumentWriter(output)
    story = fitz.Story(html=html, user_css=css)
    more = 1
    safety = 0
    while more and safety < 500:
        dev = writer.begin_page(mediabox)
        more, _ = story.place(content_rect)
        story.draw(dev)
        writer.end_page()
        safety += 1
    writer.close()
    output.seek(0)

    name = (request.form.get("file_name") or "document").strip() + ".pdf"
    return send_file(output, mimetype="application/pdf",
                     as_attachment=True, download_name=name)


@bp.route("/md-to-docx", methods=["POST"])
def md_to_docx():
    """Markdown → .docx by walking an HTML tree built from Markdown."""
    import markdown as md_lib
    import re
    from html.parser import HTMLParser
    from docx.shared import Pt, RGBColor

    md_text = request.form.get("markdown", "").strip()
    if not md_text:
        return jsonify(error="Please enter some Markdown."), 400

    html = md_lib.markdown(
        md_text,
        extensions=["extra", "sane_lists", "fenced_code", "tables"],
    )

    docx = DocxDocument()

    class MdHTMLParser(HTMLParser):
        def __init__(self):
            super().__init__()
            self.stack: list[str] = []
            self.current_para = None
            self.list_stack: list[str] = []  # "ul" or "ol"
            self.in_pre = False
            self.pending_href: str | None = None
            self._run_formats: list[dict] = []

        def _new_paragraph(self, style=None):
            self.current_para = docx.add_paragraph(style=style) if style else docx.add_paragraph()
            return self.current_para

        def _add_run(self, text):
            if not text:
                return
            p = self.current_para or self._new_paragraph()
            run = p.add_run(text)
            fmt = {}
            for f in self._run_formats:
                fmt.update(f)
            if fmt.get("bold"): run.bold = True
            if fmt.get("italic"): run.italic = True
            if fmt.get("code") or self.in_pre:
                run.font.name = "Consolas"
                run.font.size = Pt(10)
            if self.pending_href:
                run.font.color.rgb = RGBColor(0x1A, 0x0D, 0xAB)
                run.underline = True

        def handle_starttag(self, tag, attrs):
            self.stack.append(tag)
            if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = min(int(tag[1]), 4)
                self._new_paragraph(style=f"Heading {level}")
            elif tag == "p":
                if not (self.list_stack or self.in_pre):
                    self._new_paragraph()
            elif tag in ("ul", "ol"):
                self.list_stack.append(tag)
            elif tag == "li":
                style = "List Number" if self.list_stack and self.list_stack[-1] == "ol" else "List Bullet"
                try:
                    self._new_paragraph(style=style)
                except KeyError:
                    self._new_paragraph()
            elif tag in ("strong", "b"):
                self._run_formats.append({"bold": True})
            elif tag in ("em", "i"):
                self._run_formats.append({"italic": True})
            elif tag == "code":
                self._run_formats.append({"code": True})
            elif tag == "pre":
                self.in_pre = True
                self._new_paragraph()
            elif tag == "blockquote":
                try:
                    self._new_paragraph(style="Intense Quote")
                except KeyError:
                    self._new_paragraph()
            elif tag == "a":
                href = dict(attrs).get("href", "")
                self.pending_href = href
            elif tag == "hr":
                docx.add_paragraph("─" * 40)
            elif tag == "br":
                if self.current_para is not None:
                    self.current_para.add_run().add_break()

        def handle_endtag(self, tag):
            if self.stack and self.stack[-1] == tag:
                self.stack.pop()
            if tag in ("ul", "ol") and self.list_stack:
                self.list_stack.pop()
            elif tag in ("strong", "b", "em", "i", "code"):
                if self._run_formats:
                    self._run_formats.pop()
            elif tag == "pre":
                self.in_pre = False
            elif tag == "a":
                if self.pending_href:
                    self._add_run(f" ({self.pending_href})")
                self.pending_href = None

        def handle_data(self, data):
            if not data:
                return
            if self.in_pre:
                for line in data.splitlines():
                    if self.current_para is None:
                        self._new_paragraph()
                    r = self.current_para.add_run(line)
                    r.font.name = "Consolas"
                    r.font.size = Pt(10)
                    self.current_para.add_run().add_break()
            else:
                # Collapse whitespace like HTML does
                text = re.sub(r"\s+", " ", data)
                self._add_run(text)

    parser = MdHTMLParser()
    parser.feed(html)

    output = io.BytesIO()
    docx.save(output)
    output.seek(0)

    name = (request.form.get("file_name") or "document").strip() + ".docx"
    return send_file(
        output,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        as_attachment=True, download_name=name,
    )


@bp.route("/html-to-pdf", methods=["POST"])
def html_to_pdf():
    html = request.form.get("text", "").strip()
    if not html:
        return jsonify(error="Please enter some HTML content."), 400

    # Wrap in basic structure if no <html> tag present
    if "<html" not in html.lower():
        html = f"<!DOCTYPE html><html><body>{html}</body></html>"

    # Prefer LibreOffice for proper CSS rendering
    pdf_bytes = _soffice_convert(html.encode("utf-8"), "html", "pdf")

    if pdf_bytes is None:
        # Fallback: PyMuPDF's minimal HTML rendering
        doc = fitz.open()
        try:
            page = doc.new_page(width=595, height=842)  # A4
            try:
                page.insert_htmlbox(fitz.Rect(50, 50, 545, 792), html)
            except Exception as e:
                log_error(e, "html-to-pdf insert_htmlbox")
                return jsonify(error="HTML rendering failed (the markup may be invalid or use unsupported features)."), 400

            output = io.BytesIO()
            doc.save(output)
            output.seek(0)
            pdf_bytes = output.getvalue()
        finally:
            doc.close()

    output = io.BytesIO(pdf_bytes)
    output.seek(0)

    return send_file(output, mimetype="application/pdf",
                     as_attachment=True, download_name="converted.pdf")


@bp.route("/ocr-pdf", methods=["POST"])
def ocr_pdf():
    if not HAS_TESSERACT:
        return jsonify(error="OCR requires 'pytesseract' and the Tesseract binary. Install: pip install pytesseract, plus Tesseract from https://github.com/tesseract-ocr/tesseract"), 400

    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_SINGLE), 400

    mode = request.form.get("mode", "searchable")
    lang = request.form.get("lang", "eng")
    dpi = safe_int(request.form.get("dpi"), 200, min_val=72, max_val=400)

    pdf_data = files[0].read()
    try:
        src = fitz.open(stream=pdf_data, filetype="pdf")
    except Exception as e:
        log_error(e, "ocr-pdf open")
        return jsonify(error="Could not open PDF (the file may be corrupted or password-protected)."), 400

    zoom = dpi / 72

    try:
        if mode == "text":
            text_parts = []
            for i, page in enumerate(src):
                pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                with Image.open(io.BytesIO(pix.tobytes("png"))) as img:
                    text = pytesseract.image_to_string(img, lang=lang)
                text_parts.append(f"--- Page {i + 1} ---\n{text.strip()}")
            combined = "\n\n".join(text_parts).strip()
            return jsonify(text=combined or "(No text detected)")

        output = fitz.open()
        try:
            for page in src:
                pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                with Image.open(io.BytesIO(pix.tobytes("png"))) as img:
                    page_pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                        img, extension="pdf", lang=lang)
                with fitz.open(stream=page_pdf_bytes, filetype="pdf") as sub:
                    output.insert_pdf(sub)

            buf = io.BytesIO()
            output.save(buf)
            buf.seek(0)
        finally:
            output.close()

        name = files[0].filename.rsplit(".", 1)[0] + "_ocr.pdf"
        return send_file(buf, mimetype="application/pdf",
                         as_attachment=True, download_name=name)
    except pytesseract.TesseractNotFoundError:
        return jsonify(error="Tesseract binary not found. Install from https://github.com/tesseract-ocr/tesseract and ensure it is on PATH."), 400
    except Exception as e:
        msg = str(e)
        log_error(e, f"ocr-pdf lang={lang}")
        if "language" in msg.lower() or "traineddata" in msg.lower():
            return jsonify(error=f"Language pack '{lang}' not installed. Download its .traineddata file into your Tesseract tessdata directory."), 400
        return jsonify(error="OCR failed (the PDF may be image-only or unreadable)."), 400
    finally:
        src.close()


@bp.route("/cad-to-pdf", methods=["POST"])
def cad_to_pdf():
    if not HAS_EZDXF:
        return jsonify(error="CAD conversion requires 'ezdxf' and 'matplotlib'. Install: pip install ezdxf matplotlib"), 400

    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_SINGLE), 400

    target = request.form.get("format", "pdf")
    dpi = safe_int(request.form.get("dpi"), 150, min_val=72, max_val=600)

    filename = files[0].filename
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    file_data = files[0].read()

    import tempfile, os, subprocess
    with tempfile.TemporaryDirectory() as tmpdir:
        if ext == "dwg":
            if not ODA_CONVERTER:
                return jsonify(error="DWG support requires ODA File Converter. Download it free from https://www.opendesign.com/guestfiles/oda_file_converter and ensure it is on your PATH. Or convert your DWG to DXF first."), 400

            in_dir = os.path.join(tmpdir, "in")
            out_dir = os.path.join(tmpdir, "out")
            os.makedirs(in_dir)
            os.makedirs(out_dir)
            dwg_path = os.path.join(in_dir, "input.dwg")
            with open(dwg_path, "wb") as f:
                f.write(file_data)

            try:
                subprocess.run(
                    [ODA_CONVERTER, in_dir, out_dir, "ACAD2018", "DXF", "0", "1", "*.DWG"],
                    check=True, capture_output=True, timeout=60,
                )
            except subprocess.CalledProcessError as e:
                log_error(e, "cad-to-pdf ODA")
                return jsonify(error="DWG to DXF conversion failed (file may be corrupted or use an unsupported version)."), 400
            except subprocess.TimeoutExpired:
                return jsonify(error="DWG conversion timed out."), 400

            dxf_path = os.path.join(out_dir, "input.dxf")
            if not os.path.exists(dxf_path):
                return jsonify(error="DWG to DXF conversion produced no output."), 400
            doc = ezdxf.readfile(dxf_path)
        elif ext == "dxf":
            dxf_path = os.path.join(tmpdir, "input.dxf")
            with open(dxf_path, "wb") as f:
                f.write(file_data)
            try:
                doc = ezdxf.readfile(dxf_path)
            except Exception as e:
                log_error(e, "cad-to-pdf dxf parse")
                return jsonify(error="Invalid DXF file (the file may be corrupted or use an unsupported feature)."), 400
        else:
            return jsonify(error="Upload a .dxf or .dwg file."), 400

        msp = doc.modelspace()
        fig = plt.figure()
        ax = fig.add_axes([0, 0, 1, 1])
        ax.set_aspect("equal")
        ax.set_axis_off()

        try:
            ctx = RenderContext(doc)
            backend = MatplotlibBackend(ax)
            Frontend(ctx, backend).draw_layout(msp, finalize=True)
        except Exception as e:
            plt.close(fig)
            log_error(e, "cad-to-pdf render")
            return jsonify(error="CAD rendering failed (the drawing may use unsupported entities)."), 400

        buf = io.BytesIO()
        base_name = filename.rsplit(".", 1)[0]
        if target == "pdf":
            fig.savefig(buf, format="pdf", bbox_inches="tight", pad_inches=0.2)
            plt.close(fig)
            buf.seek(0)
            return send_file(buf, mimetype="application/pdf",
                             as_attachment=True, download_name=base_name + ".pdf")
        else:
            fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", pad_inches=0.2)
            plt.close(fig)
            buf.seek(0)
            return send_file(buf, mimetype="image/png",
                             as_attachment=True, download_name=base_name + ".png")


# ── PDF → PowerPoint ─────────────────────────────────────

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
SLIDE_SIZES_EMU = {
    # python-pptx uses English Metric Units (914400 EMU per inch).
    "16:9": (12192000, 6858000),   # 13.333 × 7.5 in (default widescreen)
    "4:3":  (9144000, 6858000),    # 10 × 7.5 in
    "a4":   (10692000, 7560000),   # ~11.69 × 8.27 in (landscape A4)
}


@bp.route("/pdf-to-pptx")
def pdf_to_pptx_page():
    return render_template("upload_tool.html",
        title="PDF to PowerPoint",
        description="Convert each PDF page into a slide image in a .pptx file",
        endpoint="/convert/pdf-to-pptx",
        accept=".pdf",
        multiple=False,
        options=[
            {"type": "select", "name": "slide_size", "label": "Slide size", "default": "16:9",
             "choices": [
                 {"value": "16:9", "label": "Widescreen 16:9 (default)"},
                 {"value": "4:3",  "label": "Standard 4:3"},
                 {"value": "a4",   "label": "A4 landscape"},
             ]},
            {"type": "number", "name": "dpi", "label": "Render DPI",
             "default": 150, "min": 72, "max": 300},
            {"type": "text", "name": "pages", "label": "Pages (blank = all)",
             "placeholder": "e.g. 1-3, 5, 8-10"},
        ],
        button_text="Convert to PPTX")


@bp.route("/pdf-to-pptx", methods=["POST"])
def pdf_to_pptx():
    from routes._helpers import safe_int, log_error, NO_FILE_SINGLE
    from routes.pdf_tools import parse_page_ranges

    if not HAS_PPTX:
        return jsonify(error="python-pptx is not installed. Run: pip install python-pptx"), 400

    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_SINGLE), 400

    dpi = safe_int(request.form.get("dpi"), 150, min_val=72, max_val=300)
    slide_size = request.form.get("slide_size", "16:9")
    if slide_size not in SLIDE_SIZES_EMU:
        slide_size = "16:9"
    pages_spec = (request.form.get("pages") or "").strip()

    pdf_data = files[0].read()
    try:
        doc = fitz.open(stream=pdf_data, filetype="pdf")
    except Exception as e:
        log_error(e, "pdf-to-pptx open")
        return jsonify(error="Could not open PDF (the file may be corrupted or password-protected)."), 400

    try:
        try:
            target_pages = parse_page_ranges(pages_spec, len(doc))
        except (ValueError, IndexError):
            return jsonify(error="Invalid page range. Use e.g. '1-3, 5, 8-10'."), 400

        if not target_pages:
            return jsonify(error="No pages selected."), 400

        prs = Presentation()
        slide_w, slide_h = SLIDE_SIZES_EMU[slide_size]
        prs.slide_width = slide_w
        prs.slide_height = slide_h
        blank_layout = prs.slide_layouts[6]  # 'Blank'

        mat = fitz.Matrix(dpi / 72, dpi / 72)
        for idx in target_pages:
            page = doc[idx]
            pix = page.get_pixmap(matrix=mat, alpha=False)
            png_bytes = pix.tobytes("png")
            img_w, img_h = pix.width, pix.height

            # Aspect-fit: scale to slide while preserving aspect ratio, then center.
            slide_ratio = slide_w / slide_h
            img_ratio = img_w / img_h
            if img_ratio > slide_ratio:
                draw_w = slide_w
                draw_h = int(slide_w / img_ratio)
            else:
                draw_h = slide_h
                draw_w = int(slide_h * img_ratio)
            left = (slide_w - draw_w) // 2
            top = (slide_h - draw_h) // 2

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(io.BytesIO(png_bytes),
                                     Emu(left), Emu(top),
                                     width=Emu(draw_w), height=Emu(draw_h))

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
    finally:
        doc.close()

    base = files[0].filename.rsplit(".", 1)[0]
    return send_file(output, mimetype=PPTX_MIME,
                     as_attachment=True, download_name=f"{base}.pptx")


# ── PowerPoint → PDF ─────────────────────────────────────

@bp.route("/pptx-to-pdf")
def pptx_to_pdf_page():
    return render_template("upload_tool.html",
        title="PowerPoint to PDF",
        description="Convert PowerPoint or OpenDocument presentations to PDF",
        notes=_soffice_available_notes(),
        endpoint="/convert/pptx-to-pdf",
        accept=".pptx,.ppt,.odp",
        multiple=False,
        options=[],
        button_text="Convert to PDF")


@bp.route("/pptx-to-pdf", methods=["POST"])
def pptx_to_pdf():
    import os
    import tempfile
    import subprocess
    from routes._helpers import log_error, NO_FILE_SINGLE

    if not SOFFICE:
        return jsonify(error="LibreOffice (soffice) is not installed or not on PATH. Install LibreOffice and restart the server."), 400

    files = request.files.getlist("files")
    if not files or not files[0].filename:
        return jsonify(error=NO_FILE_SINGLE), 400

    f = files[0]
    ext = f.filename.rsplit(".", 1)[-1].lower() if "." in f.filename else ""
    if ext not in ("pptx", "ppt", "odp"):
        return jsonify(error="Unsupported format. Upload .pptx, .ppt, or .odp."), 400

    safe_name = "input." + ext
    with tempfile.TemporaryDirectory() as tmp:
        in_path = os.path.join(tmp, safe_name)
        f.save(in_path)

        try:
            proc = subprocess.run(
                [SOFFICE, "--headless", "--convert-to", "pdf",
                 "--outdir", tmp, in_path],
                capture_output=True, timeout=300,
            )
        except subprocess.TimeoutExpired:
            return jsonify(error="Conversion timed out (file is too complex or too large)."), 400
        except Exception as e:
            log_error(e, "pptx-to-pdf subprocess")
            return jsonify(error="LibreOffice failed to launch."), 400

        if proc.returncode != 0:
            err = proc.stderr.decode("utf-8", errors="replace")[:200] or "unknown error"
            log_error(RuntimeError(err), "pptx-to-pdf")
            return jsonify(error="LibreOffice could not convert this file."), 400

        out_pdf = os.path.join(tmp, "input.pdf")
        if not os.path.exists(out_pdf):
            return jsonify(error="Conversion produced no output (file may be corrupted or empty)."), 400

        with open(out_pdf, "rb") as fp:
            data = fp.read()

    base = f.filename.rsplit(".", 1)[0]
    return send_file(io.BytesIO(data), mimetype="application/pdf",
                     as_attachment=True, download_name=f"{base}.pdf")
